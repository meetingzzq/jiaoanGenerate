"""
文档处理器 - 支持提取Word、PPT、Excel、TXT等文档内容
"""
import os
import io
import zipfile
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT


def detect_file_format(file_path):
    """
    检测文件的真实格式（通过文件头）
    返回: 'docx', 'doc', 'pdf', 'zip', 'rtf', 'txt', 'unknown'
    """
    try:
        with open(file_path, 'rb') as f:
            header = f.read(8)
            
        # DOCX/XLSX/PPTX 都是 ZIP 格式，以 PK 开头
        if header[:2] == b'PK':
            return 'zip_based'
        
        # DOC (OLE格式) 以 D0 CF 11 E0 开头
        if header[:4] == b'\xD0\xCF\x11\xE0':
            return 'doc'
        
        # PDF 以 %PDF 开头
        if header[:4] == b'%PDF':
            return 'pdf'
        
        # RTF 以 {\rtf 开头
        if header[:5] == b'{\\rtf':
            return 'rtf'
        
        return 'unknown'
    except Exception as e:
        print(f"检测文件格式失败: {e}")
        return 'unknown'


def extract_text_from_docx(file_path):
    """
    从Word文档(.docx)中提取文本内容
    """
    try:
        print(f"开始读取Word文档: {file_path}")
        print(f"文件是否存在: {os.path.exists(file_path)}")
        
        if not os.path.exists(file_path):
            print(f"错误: 文件不存在 - {file_path}")
            return None
        
        file_size = os.path.getsize(file_path)
        print(f"文件大小: {file_size} 字节")
        
        if file_size < 100:
            print(f"错误: 文件太小，可能已损坏")
            return None
        
        real_format = detect_file_format(file_path)
        print(f"检测到文件真实格式: {real_format}")
        
        if real_format == 'doc':
            print("警告: 文件扩展名是.docx，但实际是旧版.doc格式(OLE)")
            print("尝试使用其他方法读取...")
            return extract_text_from_doc_ole(file_path)
        
        if real_format == 'rtf':
            print("警告: 文件扩展名是.docx，但实际是RTF格式")
            return extract_text_from_rtf(file_path)
        
        if real_format != 'zip_based':
            print(f"警告: 文件可能不是有效的docx格式 (检测为: {real_format})")
            if real_format == 'unknown':
                with open(file_path, 'rb') as f:
                    sample = f.read(100)
                print(f"文件头(十六进制): {sample[:20].hex()}")
                print(f"文件头(文本): {sample[:20]}")
        
        with open(file_path, 'rb') as f:
            file_content = f.read()
        
        if not zipfile.is_zipfile(io.BytesIO(file_content)):
            print("错误: 文件不是有效的ZIP格式，可能已损坏或格式不正确")
            return try_read_as_text(file_path)
        
        doc = Document(io.BytesIO(file_content))
        full_text = []
        
        print(f"  提取段落: {len(doc.paragraphs)} 个")
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                full_text.append(para.text)
        
        print(f"  提取表格: {len(doc.tables)} 个")
        for table_idx, table in enumerate(doc.tables, 1):
            full_text.append(f"\n--- 表格 {table_idx} ---")
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    full_text.append(' | '.join(row_text))
        
        result = '\n'.join(full_text)
        print(f"  提取完成，总字符数: {len(result)}")
        return result if result.strip() else None
    except Exception as e:
        print(f"读取Word文档失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return try_read_as_text(file_path)


def try_read_as_text(file_path):
    """
    尝试将文件作为纯文本读取（最后的回退方案）
    """
    try:
        print("尝试作为纯文本读取...")
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                    if content and len(content) > 10:
                        print(f"  使用编码 {encoding} 成功读取，字符数: {len(content)}")
                        return content
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        print("无法作为文本读取")
        return None
    except Exception as e:
        print(f"文本读取失败: {e}")
        return None


def extract_text_from_doc_ole(file_path):
    """
    从OLE格式的.doc文件中提取文本（旧版Word格式）
    """
    try:
        print("尝试读取OLE格式的.doc文件...")
        
        try:
            import subprocess
            result = subprocess.run(['antiword', file_path], capture_output=True, text=True, timeout=30)
            if result.returncode == 0 and result.stdout.strip():
                print("  使用antiword读取成功")
                return result.stdout
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
        
        try:
            import olefile
            ole = olefile.OleFileIO(file_path)
            if ole.exists('WordDocument'):
                stream = ole.openstream('WordDocument')
                content = stream.read()
                text_parts = []
                for byte in content:
                    if 32 <= byte < 127:
                        text_parts.append(chr(byte))
                    elif byte > 127:
                        try:
                            text_parts.append(chr(byte))
                        except:
                            pass
                result = ''.join(text_parts)
                if len(result) > 50:
                    print(f"  使用olefile提取到 {len(result)} 字符")
                    return result
        except ImportError:
            print("  olefile未安装")
        except Exception as e:
            print(f"  olefile读取失败: {e}")
        
        return try_read_as_text(file_path)
        
    except Exception as e:
        print(f"读取OLE文档失败: {str(e)}")
        return None


def extract_text_from_rtf(file_path):
    """
    从RTF文件中提取文本
    """
    try:
        print("尝试读取RTF格式文件...")
        
        try:
            from striprtf.striprtf import rtf_to_text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            text = rtf_to_text(rtf_content)
            if text.strip():
                print(f"  RTF提取成功，字符数: {len(text)}")
                return text
        except ImportError:
            print("  striprtf未安装，尝试手动解析")
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        import re
        text = re.sub(r'\\[a-z]+\d*\s*', '', content)
        text = re.sub(r'[{}]', '', text)
        text = re.sub(r'\\[^a-z]', '', text)
        
        if text.strip():
            print(f"  手动RTF提取，字符数: {len(text)}")
            return text.strip()
        
        return None
    except Exception as e:
        print(f"读取RTF失败: {e}")
        return None


def extract_text_from_doc(file_path):
    """
    从旧版Word文档(.doc)中提取文本内容
    使用antiword或textract等工具
    """
    try:
        # 尝试使用antiword
        import subprocess
        result = subprocess.run(['antiword', file_path], capture_output=True, text=True)
        if result.returncode == 0:
            return result.stdout
        else:
            print(f"antiword读取失败: {result.stderr}")
            return None
    except FileNotFoundError:
        print("antiword未安装，无法读取.doc文件")
        return None
    except Exception as e:
        print(f"读取旧版Word文档失败: {str(e)}")
        return None


def extract_text_from_pptx(file_path):
    """
    从PowerPoint文档(.pptx)中提取文本内容
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = Presentation(file_path)
        full_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            full_text.append(f"\n--- 第{slide_num}页 ---\n")
            slide_texts = []
            
            for shape in slide.shapes:
                # 提取文本框和占位符的文本
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                
                # 提取表格内容
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        table = shape.table
                        table_texts = []
                        for row in table.rows:
                            row_texts = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_texts.append(cell.text.strip())
                            if row_texts:
                                table_texts.append(' | '.join(row_texts))
                        if table_texts:
                            slide_texts.append('\n'.join(table_texts))
                    except Exception as e:
                        print(f"  读取表格失败: {e}")
                        continue
                
                # 提取组合形状中的文本
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                slide_texts.append(sub_shape.text.strip())
                    except Exception as e:
                        print(f"  读取组合形状失败: {e}")
                        continue
            
            # 去重并添加页面文本
            seen = set()
            unique_texts = []
            for text in slide_texts:
                if text not in seen:
                    seen.add(text)
                    unique_texts.append(text)
            
            if unique_texts:
                full_text.append('\n'.join(unique_texts))
        
        result = '\n'.join(full_text)
        return result if result.strip() else None
        
    except ImportError:
        print("python-pptx未安装，无法读取PPT文件")
        return None
    except Exception as e:
        print(f"读取PPT文档失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return None


def extract_text_from_ppt(file_path):
    """
    从旧版PowerPoint文档(.ppt)中提取文本内容
    """
    try:
        # 尝试使用LibreOffice转换为pptx再读取
        import subprocess
        import tempfile
        
        with tempfile.TemporaryDirectory() as tmpdir:
            # 使用LibreOffice转换
            result = subprocess.run([
                'soffice', '--headless', '--convert-to', 'pptx',
                '--outdir', tmpdir, file_path
            ], capture_output=True, text=True)
            
            if result.returncode == 0:
                # 获取转换后的文件路径
                base_name = os.path.basename(file_path)
                new_name = os.path.splitext(base_name)[0] + '.pptx'
                new_path = os.path.join(tmpdir, new_name)
                
                if os.path.exists(new_path):
                    return extract_text_from_pptx(new_path)
            
            print(f"LibreOffice转换失败: {result.stderr}")
            return None
    except FileNotFoundError:
        print("LibreOffice未安装，无法读取.ppt文件")
        return None
    except Exception as e:
        print(f"读取旧版PPT文档失败: {str(e)}")
        return None


def extract_text_from_excel(file_path):
    """
    从Excel文档(.xlsx, .xls)中提取文本内容
    """
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        full_text = []
        
        for sheet_name in wb.sheetnames:
            full_text.append(f"\n--- 工作表: {sheet_name} ---\n")
            sheet = wb[sheet_name]
            
            # 检查工作表是否有内容
            has_content = False
            for row in sheet.iter_rows():
                row_text = []
                for cell in row:
                    if cell.value is not None:
                        row_text.append(str(cell.value))
                        has_content = True
                if row_text:
                    full_text.append(' | '.join(row_text))
            
            # 如果工作表没有内容，添加提示
            if not has_content:
                full_text.append("(空工作表)")
            
            # 添加工作表之间的分隔
            full_text.append("")
        
        result = '\n'.join(full_text)
        return result if result.strip() else None
    except ImportError:
        print("openpyxl未安装，无法读取Excel文件")
        return None
    except Exception as e:
        print(f"读取Excel文档失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return None


def extract_text_from_txt(file_path):
    """
    从文本文件(.txt)中提取内容
    """
    try:
        # 尝试多种编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                    print(f"  使用编码: {encoding}, 字符数: {len(content)}")
                    return content
            except UnicodeDecodeError:
                continue
        
        print(f"无法识别文件编码: {file_path}")
        return None
    except Exception as e:
        print(f"读取文本文件失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return None


def extract_text_from_pdf(file_path):
    """
    从PDF文档中提取文本内容
    """
    try:
        import PyPDF2
        print(f"  开始读取PDF...")
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            full_text = []
            total_pages = len(pdf_reader.pages)
            print(f"  PDF总页数: {total_pages}")
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    full_text.append(f"\n--- 第{page_num}页 ---\n")
                    full_text.append(text)
                
                # 每10页打印一次进度
                if page_num % 10 == 0:
                    print(f"  已处理 {page_num}/{total_pages} 页")
            
            result = '\n'.join(full_text)
            print(f"  PDF提取完成，总字符数: {len(result)}")
            return result if result.strip() else None
    except ImportError:
        print("PyPDF2未安装，无法读取PDF文件")
        return None
    except Exception as e:
        print(f"读取PDF文档失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return None


def extract_document_content(file_path):
    """
    根据文件类型提取文档内容
    
    Args:
        file_path: 文档文件路径
    
    Returns:
        str: 提取的文本内容，失败返回None
    """
    if not os.path.exists(file_path):
        print(f"文件不存在: {file_path}")
        return None
    
    # 获取文件扩展名
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    # 根据文件类型选择提取方法
    extractors = {
        '.docx': extract_text_from_docx,
        '.doc': extract_text_from_doc,
        '.pptx': extract_text_from_pptx,
        '.ppt': extract_text_from_ppt,
        '.xlsx': extract_text_from_excel,
        '.xls': extract_text_from_excel,
        '.txt': extract_text_from_txt,
        '.pdf': extract_text_from_pdf,
    }
    
    if ext in extractors:
        return extractors[ext](file_path)
    else:
        print(f"不支持的文件格式: {ext}")
        return None


def get_document_summary(file_path, max_length=500):
    """
    获取文档摘要（前max_length个字符）
    
    Args:
        file_path: 文档文件路径
        max_length: 摘要最大长度
    
    Returns:
        str: 文档摘要
    """
    content = extract_document_content(file_path)
    if content:
        # 去除多余空白并截取前max_length个字符
        content = ' '.join(content.split())
        if len(content) > max_length:
            return content[:max_length] + '...'
        return content
    return None


if __name__ == '__main__':
    # 测试代码
    import sys
    
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        print(f"测试文件: {test_file}")
        print("=" * 50)
        
        content = extract_document_content(test_file)
        if content:
            # 打印全部内容
            print(content)
            print("\n" + "=" * 50)
            print(f"总字符数: {len(content)}")
            
            # 同时保存到文件方便查看
            output_file = test_file + '.extracted.txt'
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(content)
            print(f"完整内容已保存到: {output_file}")
        else:
            print("无法提取内容")
    else:
        print("用法: python document_processor.py <文件路径>")
