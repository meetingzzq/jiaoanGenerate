"""
详细测试PPT提取，对比每一页的内容
"""
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

file_path = '../../test_documents/test_documents/test_ppt.pptx'

print(f"测试文件: {file_path}")
print("=" * 80)

prs = Presentation(file_path)

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n{'='*80}")
    print(f"第 {slide_num} 页")
    print(f"{'='*80}")
    
    shape_count = 0
    for shape in slide.shapes:
        shape_count += 1
        print(f"\n  Shape {shape_count}:")
        print(f"    类型: {shape.shape_type}")
        print(f"    名称: {shape.name}")
        
        # 检查是否有文本
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                print(f"    文本内容:")
                print(f"      {text[:200]}{'...' if len(text) > 200 else ''}")
            else:
                print(f"    文本: (空)")
        
        # 检查是否是表格
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            print(f"    这是一个表格!")
            try:
                table = shape.table
                print(f"    表格行数: {len(table.rows)}")
                print(f"    表格列数: {len(table.columns)}")
            except Exception as e:
                print(f"    读取表格失败: {e}")
        
        # 检查是否是组合
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"    这是一个组合形状!")
            try:
                sub_count = 0
                for sub_shape in shape.shapes:
                    sub_count += 1
                    if hasattr(sub_shape, "text") and sub_shape.text.strip():
                        print(f"      子Shape {sub_count} 文本: {sub_shape.text.strip()[:100]}")
            except Exception as e:
                print(f"    读取组合失败: {e}")

print(f"\n{'='*80}")
print("测试完成")
